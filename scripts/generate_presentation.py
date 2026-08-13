#!/usr/bin/env python3
"""Generate Allo Inventory Reservations project presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BRAND = RGBColor(0x0F, 0x76, 0x6E)       # teal
BRAND_DARK = RGBColor(0x11, 0x5E, 0x59)
INK = RGBColor(0x18, 0x21, 0x2F)
MUTED = RGBColor(0x64, 0x70, 0x84)
DANGER = RGBColor(0xB4, 0x23, 0x18)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xEE, 0xF3, 0xF7)
ACCENT_LIGHT = RGBColor(0xCC, 0xFB, 0xF1)
SUCCESS = RGBColor(0x05, 0x96, 0x69)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide, color=LIGHT_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    sp = slide.shapes._spTree
    sp.remove(bg._element)
    sp.insert(2, bg._element)


def add_accent_bar(slide, height=Inches(0.08)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()


def add_footer(slide, text="Allo Inventory Reservations"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT


def add_title_slide(title, subtitle, tagline=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BRAND)
    # decorative circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x0D, 0x6B, 0x63)
    circle.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10), Inches(1.5))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(10), Inches(1))
    tf = sbox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xCC, 0xFB, 0xF1)

    if tagline:
        tbox2 = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(10), Inches(0.6))
        tf2 = tbox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = tagline
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0x99, 0xF6, 0xE4)


def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BRAND_DARK)
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11), Inches(1.5))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_content_slide(title, bullets, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide)

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.8))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BRAND_DARK

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12), Inches(0.5))
        stf = sbox.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = MUTED

    y_start = 1.6 if subtitle else 1.3
    bbox = slide.shapes.add_textbox(Inches(0.6), Inches(y_start), Inches(12), Inches(5.2))
    btf = bbox.text_frame
    btf.word_wrap = True
    for i, item in enumerate(bullets):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = item
            p.level = 0
        p.font.size = Pt(16 if p.level == 0 else 14)
        p.font.color.rgb = INK if p.level == 0 else MUTED
        p.space_after = Pt(8)


def add_two_column_slide(title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide)

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.8))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BRAND_DARK

    for col, (ctitle, citems, x) in enumerate([
        (left_title, left_items, 0.6),
        (right_title, right_items, 6.8),
    ]):
        hbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(5.8), Inches(5.5))
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = WHITE
        hbox.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

        ht = slide.shapes.add_textbox(Inches(x + 0.25), Inches(1.45), Inches(5.3), Inches(0.5))
        htf = ht.text_frame
        hp = htf.paragraphs[0]
        hp.text = ctitle
        hp.font.size = Pt(18)
        hp.font.bold = True
        hp.font.color.rgb = BRAND

        cbox = slide.shapes.add_textbox(Inches(x + 0.25), Inches(2.0), Inches(5.3), Inches(4.5))
        ctf = cbox.text_frame
        ctf.word_wrap = True
        for i, item in enumerate(citems):
            if i == 0:
                cp = ctf.paragraphs[0]
            else:
                cp = ctf.add_paragraph()
            cp.text = f"• {item}"
            cp.font.size = Pt(14)
            cp.font.color.rgb = INK
            cp.space_after = Pt(6)


def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide)

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.8))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BRAND_DARK

    cols = len(headers)
    nrows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(nrows, cols, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4 * nrows))
    table = tbl_shape.table

    col_width = Inches(12 / cols)
    for c in range(cols):
        table.columns[c].width = col_width

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(12)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.color.rgb = INK


def add_flow_slide(title, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide)

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.8))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BRAND_DARK

    n = len(steps)
    box_w = min(2.0, 11.5 / n)
    gap = (11.5 - box_w * n) / max(n - 1, 1)
    x_start = 0.6

    for i, (step_title, step_desc, color) in enumerate(steps):
        x = x_start + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(box_w), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        st = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.15), Inches(box_w - 0.2), Inches(0.5))
        stf = st.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = step_title
        sp.font.size = Pt(13)
        sp.font.bold = True
        sp.font.color.rgb = WHITE
        sp.alignment = PP_ALIGN.CENTER

        sd = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.7), Inches(box_w - 0.2), Inches(0.9))
        sdf = sd.text_frame
        sdf.word_wrap = True
        sdp = sdf.paragraphs[0]
        sdp.text = step_desc
        sdp.font.size = Pt(10)
        sdp.font.color.rgb = RGBColor(0xE0, 0xF2, 0xF1)
        sdp.alignment = PP_ALIGN.CENTER

        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + box_w + gap * 0.15), Inches(2.7), Inches(gap * 0.7), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()


def add_closing_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BRAND)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x0D, 0x6B, 0x63)
    circle.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(1.2))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11), Inches(1))
    stf = sbox.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Questions & Discussion"
    sp.font.size = Pt(24)
    sp.font.color.rgb = RGBColor(0xCC, 0xFB, 0xF1)

    lbox = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11), Inches(0.8))
    ltf = lbox.text_frame
    lp = ltf.paragraphs[0]
    lp.text = "github.com/moganraj05/allo-inventory-reservation"
    lp.font.size = Pt(14)
    lp.font.color.rgb = RGBColor(0x99, 0xF6, 0xE4)


# ── BUILD SLIDES ──────────────────────────────────────────────────────────────

add_title_slide(
    "Allo Inventory Reservations",
    "Multi-Warehouse Stock Reservation System",
    "Preventing Overselling at Checkout  •  Next.js 15  •  PostgreSQL  •  Prisma"
)

add_content_slide("Agenda", [
    "Problem Statement & Business Context",
    "Solution Overview & Key Features",
    "Technology Stack",
    "System Architecture & Design Patterns",
    "Database Schema & Data Model",
    "Reservation Lifecycle & State Machine",
    "Concurrency & Race Condition Handling",
    "Expiry Strategy & Idempotency",
    "API Design & Endpoints",
    "Frontend UI/UX Design",
    "Deployment & DevOps",
    "Trade-offs & Future Roadmap",
])

add_section_slide("Problem & Solution")

add_content_slide(
    "The Problem: Overselling at Checkout",
    [
        "E-commerce platforms face a critical challenge when multiple customers attempt to purchase the last available unit simultaneously.",
        "Without proper inventory locking, race conditions cause overselling — selling more units than physically available.",
        "Abandoned checkouts leave inventory locked indefinitely, reducing sellable stock.",
        "Multi-warehouse operations add complexity: stock must be tracked per location with accurate availability.",
        "Payment failures must gracefully return units to available inventory without manual intervention.",
    ],
    subtitle="Why inventory reservation matters"
)

add_content_slide(
    "Our Solution",
    [
        "Temporary 10-minute stock holds when a customer proceeds to checkout.",
        "Atomic, database-level concurrency control prevents overselling under high traffic.",
        "Three-state reservation lifecycle: PENDING → CONFIRMED or RELEASED.",
        "Automatic expiry cleanup returns abandoned holds to available inventory.",
        "Idempotent API design ensures safe retries during payment processing.",
        "Multi-warehouse support with per-location stock tracking (total, reserved, available).",
    ],
    subtitle="A production-ready reservation system built for reliability"
)

add_section_slide("Technology Stack")

add_two_column_slide(
    "Technology Stack",
    "Frontend & Framework",
    [
        "Next.js 15.3 — App Router, API Routes",
        "React 19 — Client & Server Components",
        "TypeScript 5.8 — Strict type safety",
        "Custom CSS — Teal brand design system",
        "No external UI library — lightweight & fast",
    ],
    "Backend & Infrastructure",
    [
        "PostgreSQL — Hosted (Supabase / Neon)",
        "Prisma 6.8 — ORM & migrations",
        "Zod 3.24 — Request validation",
        "Vercel — Serverless deployment",
        "ESLint 9 — Code quality",
    ],
)

add_section_slide("Architecture & Design")

add_content_slide(
    "System Architecture",
    [
        "Layered Architecture — Clean separation of concerns across four tiers:",
        ("Presentation Layer — React client components (inventory-app, reservation-detail)", 1),
        ("API Layer — Next.js App Router route handlers (/api/*)", 1),
        ("Business Logic Layer — lib/reservations.ts, idempotency.ts, validation.ts", 1),
        ("Data Access Layer — Prisma ORM → PostgreSQL", 1),
        "",
        "Key Design Patterns:",
        ("Transactional boundaries — All stock mutations in prisma.$transaction() with 15s timeout", 1),
        ("Optimistic concurrency — Conditional updateMany on reservation status", 1),
        ("Idempotency middleware — DB-backed deduplication via runIdempotently()", 1),
        ("Lazy + scheduled cleanup — Dual expiry strategy for reliability", 1),
        ("Raw SQL for critical paths — Atomic stock updates via $executeRaw", 1),
    ],
    subtitle="Four-tier layered architecture with production patterns"
)

add_content_slide(
    "Project Structure",
    [
        "app/ — Next.js App Router (pages + API routes)",
        ("  page.tsx — Home / product listing", 1),
        ("  reservations/[id]/page.tsx — Checkout page", 1),
        ("  api/products, warehouses, reservations, cron — REST endpoints", 1),
        "components/ — React UI components",
        ("  inventory-app.tsx — Product grid, stock table, reserve actions", 1),
        ("  reservation-detail.tsx — Countdown timer, confirm/cancel", 1),
        "lib/ — Core business logic",
        ("  reservations.ts — Create, confirm, release, expiry cleanup", 1),
        ("  idempotency.ts — Safe retry handling", 1),
        "prisma/ — Schema, migrations, seed data",
        "scripts/ — Concurrency smoke test",
    ],
    subtitle="Organized for maintainability and clear separation"
)

add_section_slide("Database Design")

add_content_slide(
    "Database Schema — 5 Tables",
    [
        "Product — id, name, sku (unique), description",
        "Warehouse — id, name, code (unique), city",
        "StockLevel — productId + warehouseId (unique), totalUnits, reservedUnits",
        ("availableUnits = totalUnits - reservedUnits (computed, not stored)", 1),
        "Reservation — productId, warehouseId, quantity, status, expiresAt",
        ("Status enum: PENDING | CONFIRMED | RELEASED", 1),
        ("Indexes: (status, expiresAt), (productId, warehouseId)", 1),
        "IdempotencyRecord — key + method + path (unique), response JSON cache",
        "",
        "Seed Data: 5 warehouses, 10 products, 50 stock levels, 4 sample reservations",
    ],
    subtitle="PostgreSQL with Prisma ORM"
)

add_flow_slide(
    "Reservation Lifecycle",
    [
        ("RESERVE", "Customer clicks\nReserve button", RGBColor(0x25, 0x63, 0xEB)),
        ("PENDING", "10-min hold\nUnits reserved", RGBColor(0xD9, 0x77, 0x06)),
        ("CONFIRM", "Payment succeeds\nStock decremented", SUCCESS),
        ("RELEASE", "Cancel or expire\nUnits returned", DANGER),
    ]
)

add_content_slide(
    "State Transitions Detail",
    [
        "PENDING — Created on reserve. reservedUnits incremented. expiresAt = now + 10 min.",
        "CONFIRMED — Payment success. totalUnits AND reservedUnits both decremented. Permanent sale.",
        "RELEASED — Cancel or expiry. Only reservedUnits decremented. Units return to available.",
        "",
        "Guards & Validations:",
        ("Cannot confirm a RELEASED reservation → 409 RESERVATION_RELEASED", 1),
        ("Cannot confirm an expired reservation → 410 RESERVATION_EXPIRED", 1),
        ("Cannot release a CONFIRMED reservation → 409 ALREADY_CONFIRMED", 1),
        ("Insufficient stock on reserve → 409 INSUFFICIENT_STOCK", 1),
        ("Conditional status update prevents double-confirm or double-release", 1),
    ],
    subtitle="Every transition is atomic and guarded"
)

add_section_slide("Concurrency & Reliability")

add_content_slide(
    "Concurrency Handling",
    [
        "Core mechanism — Single guarded SQL UPDATE inside a Prisma transaction:",
        "",
        'UPDATE "StockLevel" SET "reservedUnits" = "reservedUnits" + qty',
        'WHERE productId = ? AND warehouseId = ?',
        '  AND (totalUnits - reservedUnits) >= qty',
        "",
        "Postgres row-level lock ensures only one transaction updates a stock row at a time.",
        "If two requests race for the last unit: first wins, second gets 0 rows updated → 409.",
        "",
        "Verified by concurrency smoke test:",
        ("8 parallel requests on 1 available unit → 1 success + 7 failures (409)", 1),
        ("Run: npm run test:concurrency", 1),
    ],
    subtitle="Race-condition safe by design"
)

add_two_column_slide(
    "Expiry & Idempotency",
    "Dual Expiry Strategy",
    [
        "Lazy Cleanup — Runs before product reads and new reservations during normal traffic",
        "Cron Endpoint — POST /api/cron/release-expired every 5 min (Vercel)",
        "Batch processing — Up to 100 expired holds per transaction",
        "Configurable window — RESERVATION_WINDOW_MINUTES (default: 10)",
        "CRON_SECRET — Bearer auth for production cron calls",
    ],
    "Idempotency Support",
    [
        "Idempotency-Key header on POST /api/reservations and confirm",
        "Stored in IdempotencyRecord table (key + method + path)",
        "Retries return cached response — no duplicate side effects",
        "In-progress duplicates → 409 IDEMPOTENCY_IN_PROGRESS",
        "Critical for payment integration (Stripe, Razorpay webhooks)",
    ],
)

add_section_slide("API & Frontend")

add_table_slide(
    "REST API Endpoints",
    ["Method", "Endpoint", "Purpose", "Auth"],
    [
        ["GET", "/api/products", "List products with per-warehouse stock", "—"],
        ["GET", "/api/warehouses", "List all warehouses", "—"],
        ["POST", "/api/reservations", "Create reservation (hold units)", "Idempotency-Key"],
        ["GET", "/api/reservations/:id", "Get reservation details", "—"],
        ["POST", "/api/reservations/:id/confirm", "Confirm pending hold", "Idempotency-Key"],
        ["POST", "/api/reservations/:id/release", "Cancel/release hold", "—"],
        ["POST", "/api/cron/release-expired", "Batch release expired holds", "Bearer CRON_SECRET"],
    ]
)

add_content_slide(
    "Error Handling Taxonomy",
    [
        "400 VALIDATION_ERROR — Invalid request body (Zod validation)",
        "401 UNAUTHORIZED — Cron endpoint without valid CRON_SECRET",
        "404 NOT_FOUND — Reservation not found",
        "409 INSUFFICIENT_STOCK — Not enough available units",
        "409 RESERVATION_RELEASED — Already released, cannot confirm",
        "409 ALREADY_CONFIRMED — Cannot release confirmed reservation",
        "409 IDEMPOTENCY_IN_PROGRESS — Duplicate request still processing",
        "410 RESERVATION_EXPIRED — Hold expired before confirmation",
        "500 INTERNAL_ERROR — Unexpected server error",
    ],
    subtitle="Structured error codes for client-side handling"
)

add_content_slide(
    "Frontend UI/UX Design",
    [
        "Design System — Custom CSS with teal brand palette (#0f766e)",
        ("Colors: Brand teal, ink (#18212f), muted (#647084), danger (#b42318)", 1),
        ("Typography: Inter font family with system-ui fallback", 1),
        ("Background: Light gray (#eef3f7) with white surface cards", 1),
        "",
        "Page 1 — Product Listing (/)",
        ("Metric row: Available / Reserved / Total counts in 3-column grid", 1),
        ("Product cards with SKU badge and per-warehouse stock table", 1),
        ("Reserve button per warehouse row (disabled when stock = 0)", 1),
        "",
        "Page 2 — Checkout (/reservations/:id)",
        ("Color-coded status badge (pending / confirmed / released)", 1),
        ("Large MM:SS countdown timer with auto-expiry handling", 1),
        ("Confirm purchase & Cancel buttons with success/cancelled popups", 1),
        ("Responsive design — Mobile breakpoint at 760px", 1),
    ],
    subtitle="Clean, modern interface matching the brand identity"
)

add_content_slide(
    "User Journey Flow",
    [
        "1. Browse Products — Customer views inventory listing with per-warehouse availability",
        "2. Reserve Unit — Clicks 'Reserve' → API creates PENDING hold → redirects to checkout",
        "3. Checkout Timer — 10-minute countdown displayed with product & warehouse details",
        "4a. Confirm Purchase — Payment succeeds → stock permanently decremented → success popup",
        "4b. Cancel Order — Customer cancels → units returned to available → cancelled popup",
        "4c. Timer Expires — Hold auto-released via lazy cleanup or cron → 410 on confirm attempt",
        "",
        "Frontend sends Idempotency-Key (crypto.randomUUID()) on reserve and confirm requests.",
        "All API fetches use cache: 'no-store' for real-time stock accuracy.",
    ],
    subtitle="End-to-end customer experience"
)

add_section_slide("Deployment & Future")

add_content_slide(
    "Deployment on Vercel",
    [
        "Platform — Vercel serverless with Next.js 15 App Router",
        "Database — Hosted PostgreSQL (Supabase or Neon) with pooled + direct connections",
        "",
        "Environment Variables:",
        ("DATABASE_URL — Pooled Postgres connection", 1),
        ("DIRECT_URL — Direct connection for Prisma migrations", 1),
        ("CRON_SECRET — Bearer token for cron endpoint auth", 1),
        ("RESERVATION_WINDOW_MINUTES — Hold duration (default: 10)", 1),
        "",
        "Deploy Steps: Push to GitHub → Import on Vercel → Set env vars → Deploy → Run migrations & seed",
        "Live Demo: allo-inventory-reservation-2.vercel.app",
    ],
    subtitle="Production-ready serverless deployment"
)

add_two_column_slide(
    "Trade-offs & Future Roadmap",
    "Current Trade-offs",
    [
        "Frontend reserves 1 unit per click (API supports up to 99)",
        "Batch expiry processes 100 reservations per transaction",
        "Postgres row locks instead of Redis (simpler, sufficient)",
        "No SQLite fallback — requires hosted Postgres",
        "Vercel free tier limits cron frequency; lazy cleanup compensates",
        "No unit/E2E tests yet (concurrency smoke test only)",
    ],
    "Planned Improvements",
    [
        "Multi-unit shopping cart (reserve 2, 5, 10 units)",
        "Admin dashboard for reservation analytics",
        "Payment integration (Razorpay, Stripe webhooks)",
        "E2E tests with Playwright",
        "Better warehouse search & filtering",
        "Background workers for large-scale expiry cleanup",
    ],
)

add_content_slide(
    "Key Takeaways",
    [
        "Prevents overselling with atomic, database-level concurrency control.",
        "Production-ready patterns: idempotency, transactional boundaries, structured errors.",
        "Dual expiry strategy ensures inventory is never permanently locked.",
        "Clean layered architecture — easy to extend with payments, admin, analytics.",
        "Modern stack: Next.js 15, React 19, Prisma, PostgreSQL, TypeScript, Vercel.",
        "Verified concurrency safety with automated smoke test.",
    ],
    subtitle="A reliable foundation for e-commerce inventory management"
)

add_closing_slide()

# Save
output_path = "/workspace/Allo_Inventory_Reservations_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
